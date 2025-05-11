"""
PDFをPPTXに変換するためのコアモジュール
PyMuPDF（fitz）を使用してPDFから画像への変換を行います
"""
import os
import re
import sys
import shutil
import tempfile
from glob import glob
from pptx import Presentation
from pptx.util import Pt
from PIL import Image
import fitz  # PyMuPDF


class PDFConverter:
    """PDFファイルをPowerPointプレゼンテーションに変換するクラス
    
    PyMuPDFライブラリを使用してPDFから画像への変換を行い、
    python-pptxライブラリを使用してPowerPointファイルを生成します。
    """
    
    def __init__(self):
        """初期化メソッド"""
        self.temp_folder = None
        self.output_folder = None
        self.image_format = "jpg"  # 画像フォーマット（jpg, png）
        self.dpi = 300  # 画像変換の解像度
    
    def convert_pdf_to_pptx(self, pdf_path, output_folder=None, callback=None):
        """
        PDFファイルをPPTXに変換する

        Args:
            pdf_path (str): 変換するPDFファイルのパス
            output_folder (str, optional): 出力先フォルダのパス。指定がなければPDFと同じ場所
            callback (callable, optional): 進捗状況を通知するコールバック関数

        Returns:
            tuple: (PPTXファイルのパス, 画像フォルダのパス)

        Raises:
            FileNotFoundError: PDFファイルが見つからない場合
            ValueError: PDFファイルでない場合や、変換中のエラー
        """
        # 入力ファイル検証
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDFファイル '{pdf_path}' が見つかりません")
        
        base, ext = os.path.splitext(pdf_path)
        if ext.lower() != '.pdf':
            raise ValueError(f"'{pdf_path}' はPDFファイルではありません")
        
        # 出力フォルダ設定
        if output_folder is None:
            output_folder = os.path.dirname(pdf_path)
        
        self.output_folder = output_folder
        
        # コールバック関数が指定されていない場合は、ダミー関数を使用
        if callback is None:
            def callback(status, message, progress=None):
                pass
        
        # PDF変換処理の開始
        callback("開始", "変換を開始します", 0)
        
        try:
            # 一時フォルダの準備
            self._setup_temp_folder(pdf_path)
            
            # PDFを画像に変換
            callback("変換中", "PDFを画像に変換しています", 10)
            image_files = self._convert_pdf_to_images(pdf_path, callback)
            
            # 画像ファイルをPowerPointスライドに配置
            callback("変換中", "PowerPointスライドを作成しています", 50)
            pptx_path = self._create_pptx_from_images(image_files, os.path.basename(pdf_path), callback)
            
            # 画像フォルダを出力先にコピー
            callback("保存中", "ファイルを保存しています", 90)
            images_folder_path = self._copy_images_to_output(os.path.basename(pdf_path))
            
            callback("完了", "変換が完了しました", 100)
            
            # 出力ファイルのパスを返す
            return pptx_path, images_folder_path
            
        except Exception as e:
            # エラーメッセージの改善
            error_msg = str(e)
            
            callback("エラー", f"変換中にエラーが発生しました: {error_msg}", None)
            raise ValueError(f"PDF変換エラー: {error_msg}") from e
        
        finally:
            # 常に一時フォルダを削除
            self._cleanup_temp_folder()
    
    def _setup_temp_folder(self, pdf_path):
        """一時作業フォルダを設定する"""
        # 一時フォルダを作成
        self.temp_folder = tempfile.mkdtemp(prefix="pdf2pptx_")
        return self.temp_folder
    
    def _cleanup_temp_folder(self):
        """一時フォルダを削除する"""
        if self.temp_folder and os.path.exists(self.temp_folder):
            shutil.rmtree(self.temp_folder, ignore_errors=True)
            self.temp_folder = None
    
    def _convert_pdf_to_images(self, pdf_path, callback):
        """
        PyMuPDFを使用してPDFを画像に変換する
        
        Args:
            pdf_path (str): 変換するPDFファイルのパス
            callback (callable): 進捗状況を通知するコールバック関数
            
        Returns:
            list: 生成された画像ファイルのパスリスト
            
        Raises:
            ValueError: PDF変換中のエラー
        """
        images_folder = os.path.join(self.temp_folder, "images")
        os.makedirs(images_folder, exist_ok=True)
        
        try:
            # PyMuPDFを使用してPDFを開く
            pdf_document = fitz.open(pdf_path)
            total_pages = len(pdf_document)
            
            image_files = []
            
            # PowerPointの最大サイズ (56インチ = 約4032ピクセル@72dpi)
            MAX_PPT_SIZE = 4032
            
            # 各ページを画像として保存
            for i in range(total_pages):
                # ページを取得
                page = pdf_document[i]
                
                # ページサイズを取得
                page_rect = page.rect
                
                # サイズに基づいて適切なDPI/ズームを計算
                # PowerPointの制限を超えないようにする
                width_pt = page_rect.width
                height_pt = page_rect.height
                
                # 最大DPIを計算（PowerPointの制限を考慮）
                max_zoom_width = MAX_PPT_SIZE / width_pt
                max_zoom_height = MAX_PPT_SIZE / height_pt
                max_zoom = min(max_zoom_width, max_zoom_height, self.dpi / 72)
                
                # 安全マージンを取る (90%)
                safe_zoom = max_zoom * 0.9
                
                # 最終的なズーム値を決定
                zoom = min(self.dpi / 72, safe_zoom)
                
                matrix = fitz.Matrix(zoom, zoom)
                
                # ページを画像としてレンダリング
                pix = page.get_pixmap(matrix=matrix)
                
                # 画像ファイルのパスを設定
                image_path = os.path.join(images_folder, f"page_{i+1:03d}.{self.image_format}")
                
                # 画像として保存
                if self.image_format.lower() == "jpg":
                    pix.save(image_path, "jpeg")
                else:
                    pix.save(image_path)
                    
                image_files.append(image_path)
                
                # 進捗状況をコールバックで通知
                progress = 10 + (i + 1) / total_pages * 40  # 10%〜50%の範囲で進捗
                callback("変換中", f"PDFを画像に変換しています ({i+1}/{total_pages})", progress)
            
            return image_files
            
        except Exception as e:
            raise ValueError(f"PDF変換エラー: {str(e)}") from e
    
    def _create_pptx_from_images(self, image_files, base_name, callback):
        """画像ファイルからPPTXを作成する"""
        if not image_files:
            raise ValueError("変換する画像ファイルがありません")
            
        # PPTXファイルのパスを設定
        pptx_filename = os.path.splitext(base_name)[0] + ".pptx"
        pptx_path = os.path.join(self.output_folder, pptx_filename)
        
        # 最初の画像からサイズを取得
        with Image.open(image_files[0]) as img:
            width, height = img.size
        
        # Presentationインスタンス作成
        prs = Presentation()
        
        # スライドサイズを画像サイズに合わせる
        prs.slide_width = Pt(width)
        prs.slide_height = Pt(height)
        
        # 白紙レイアウトを使用
        blank_layout = prs.slide_layouts[6]
        
        total_images = len(image_files)
        
        # 各画像をスライドに配置
        for i, img_path in enumerate(image_files):
            # スライド作成
            slide = prs.slides.add_slide(blank_layout)
            
            # 画像の挿入 - スライド全体に拡大表示するために左上(0,0)から開始
            pic = slide.shapes.add_picture(img_path, 0, 0)
            
            # 画像をスライド全体に拡大（アスペクト比を維持せず、完全にカバー）
            pic.width = prs.slide_width
            pic.height = prs.slide_height
            
            # 進捗状況をコールバックで通知
            progress = 50 + (i + 1) / total_images * 40  # 50%〜90%の範囲で進捗
            callback("変換中", f"PowerPointスライドを作成しています ({i+1}/{total_images})", progress)
        
        # プレゼンテーションを保存
        prs.save(pptx_path)
        
        return pptx_path
    
    def _copy_images_to_output(self, base_name):
        """変換した画像ファイルを出力フォルダにコピーする"""
        # 画像フォルダ名を設定
        images_folder_name = os.path.splitext(base_name)[0] + "_images"
        images_output_path = os.path.join(self.output_folder, images_folder_name)
        
        # 既存の画像フォルダがある場合は削除
        if os.path.exists(images_output_path):
            shutil.rmtree(images_output_path, ignore_errors=True)
        
        # 一時フォルダの画像を出力フォルダにコピー
        temp_images_folder = os.path.join(self.temp_folder, "images")
        if os.path.exists(temp_images_folder):
            shutil.copytree(temp_images_folder, images_output_path)
        
        return images_output_path


# コマンドラインから直接実行された場合
if __name__ == "__main__":
    if len(sys.argv) > 1:
        pdf_file = sys.argv[1]
        if os.path.exists(pdf_file) and pdf_file.lower().endswith('.pdf'):
            # 変換の実行
            converter = PDFConverter()
            
            def show_progress(status, message, progress=None):
                """進捗を表示するコールバック関数"""
                if progress is not None:
                    print(f"{status}: {message} - {progress:.1f}%")
                else:
                    print(f"{status}: {message}")
            
            try:
                output_pptx, output_images = converter.convert_pdf_to_pptx(
                    pdf_file, callback=show_progress
                )
                print(f"\n変換が完了しました！")
                print(f"PowerPointファイル: {output_pptx}")
                print(f"画像フォルダ: {output_images}")
            except Exception as e:
                print(f"\nエラー: {str(e)}")
        else:
            print("指定されたPDFファイルが見つからないか、PDFファイルではありません。")
            print("使用法: python pdf_converter.py <PDFファイルパス>")
    else:
        print("PDFファイルのパスを指定してください。")
        print("使用法: python pdf_converter.py <PDFファイルパス>")
