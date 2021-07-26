## ウィンドウズでコンパイルする用のコード。単品では動かない。


import sys
import os
from pptx import Presentation
from pptx.util import Pt
from glob import glob
from PIL import Image
import re
import shutil
from pdf2image import convert_from_path


# popplerへの環境変数PATHを一時的に付与 ###################################
poppler_path = os.path.join(os.getcwd(), "poppler", "bin")
os.environ["PATH"] += os.pathsep + poppler_path

# 文字列から数値だけを取り出す
def str2int(str):
    i = int(re.search('\d+', str).group())
    return i

print("====================================================")
print("PDFを、PowerPointに貼り付けます。_2021/07/24")
print("author phys-ken    Twitter:@phys_ken")
print("====================================================")
print("現在の作業フォルダは_"+os.getcwd() + "_です。")
print("必要な情報を入力したら、Enterを押してください。")
print("")
print("")
fullPath = input("PDFのフルパスを入力してください>>>")
while os.path.exists(fullPath) == False:
    print("====ファイルが存在しません。====")
    fullPath = input("PDFのフルパスを入力してください>>>")

outputdir = input("出力先のフォルダのフルパスを入力してください>>>")
while os.path.isdir(outputdir) == False:
    print("====そのフォルダには出力できません。====")
    outputdir = input("出力先のフォルダを入力してください>>>")


base, ext = os.path.splitext(fullPath)

if ext == '.pdf' or ext == '.PDF':
    print(fullPath + "の変換を準備中...")
    # figtmpfigがあれば、一旦削除
    if(os.path.isdir('figtmpfig') == True):
        shutil.rmtree('figtmpfig')
    load = convert_from_path(fullPath ,poppler_path="poppler/bin" )
    if not os.path.exists("figtmpfig"):
        # figtmpfigを作る。
        os.makedirs("figtmpfig")
    img_file = "figtmpfig/output"
    for index, jpg in enumerate(load):
        jpg_name = img_file + '_' + '{0:03d}.jpg'.format(index)
        jpg.save(jpg_name)
        print("変換中_{0:03d}.jpg".format(index))
    print("pdfを画像に変換しました。")
else:
    print("pdfファイルがありません、処理をやめます。")
    sys.exit()


# 画像ファイルの読み込み
fnms = glob('./figtmpfig/*.jpg')
print(str(len(fnms)) + "枚の画像を処理します。")

# ソート
fnms.sort(key=str2int)


im = Image.open(fnms[0])
h = im.height
w = im.width

# Presentationインスタンスの作成
ppt = Presentation()

# デフォルトはA4縦にしておく
ppt.slide_height = Pt(h)
ppt.slide_width = Pt(w)

# 幅
width = ppt.slide_width
# 高さ
height = ppt.slide_height

# レイアウト, 6番は白紙
blank_slide_layout = ppt.slide_layouts[6]

# ファイル毎にループ
for fnm in fnms:
    print(fnm + "を処理します。")
    # 白紙のスライドの追加
    slide = ppt.slides.add_slide(blank_slide_layout)

    # 画像の挿入
    pic = slide.shapes.add_picture(fnm, 0, 0)

    # 中心に移動
    pic.left = int((width - pic.width) / 2)
    pic.top = int((height - pic.height) / 2)
    print(fnm + "の処理が終わった！")

# 名前をつけて保存
ppt.save(outputdir +  '/output.pptx')
#shutil.rmtree('figtmpfig') winだとうごかない。
print("変換が終わりました！！！ 出力フォルダにoutput.pptxを作成しました。")
print("変換した画像を、" + os.getcwd() + "tmpfigtmpフォルダに入れてあります。")
print("tmpfigtmpは、次回実行時に消されます。")