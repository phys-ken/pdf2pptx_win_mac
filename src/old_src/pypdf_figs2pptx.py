from pptx import Presentation
from pptx.util import Pt
from glob import glob
from PIL import Image
import re

# 文字列から数値だけを取り出す
def str2int(str):
    i = int( re.search('\d+',str).group() )
    return i


# 画像ファイルの読み込み
fnms = glob('./fig/*.jpg')
print(str(len(fnms)) + "枚の画像を処理します。")

# ソート
fnms.sort(key=str2int)


im  = Image.open(fnms[0])
h = im.height
w = h / 1.4142

# Presentationインスタンスの作成
ppt = Presentation()

#デフォルトはA4縦にしておく
ppt.slide_height=Pt(h)
ppt.slide_width=Pt(w)

# 幅
width = ppt.slide_width
# 高さ
height = ppt.slide_height

# レイアウト, 6番は白紙
blank_slide_layout = ppt.slide_layouts[6]

# ファイル毎にループ
for fnm in fnms:
    print(fnm + "を処理します。")
    # 白紙のスライドの追加
    slide = ppt.slides.add_slide(blank_slide_layout)
    

    # 画像の挿入
    pic = slide.shapes.add_picture(fnm, 0, 0)

    # 中心に移動
    pic.left = int( ( width  - pic.width  ) / 2 )
    pic.top  = int( ( height - pic.height ) / 2 )
    print(fnm + "の処理が終わった！")

# 名前をつけて保存
ppt.save('output.pptx')