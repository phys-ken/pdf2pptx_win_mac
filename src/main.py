import sys
import os
from pptx import Presentation
from pptx.util import Pt
from glob import glob
from PIL import Image
import re
import shutil
from pathlib import Path
from pdf2image import convert_from_path


# 文字列から数値だけを取り出す
def str2int(str):
    i = int(re.search('\d+', str).group())
    return i

print("作業フォルダは_" + os.getcwd() + "_です。" )

fullPath = input("PDFのフルパスを入力してください>>>")
while os.path.exists(fullPath) == False:
    print("====ファイルが存在しません。====")
    fullPath = input("PDFのフルパスを入力してください>>>")

outputdir = input("出力先のフォルダを入力してください>>>")
while os.path.isdir(outputdir) == False:
    print("====そのフォルダには出力できません。====")
    outputdir = input("出力先のフォルダを入力してください>>>")


base, ext = os.path.splitext(fullPath)

if ext == '.pdf' or ext == '.PDF':
    print(fullPath + "の変換を準備中...")
    # figがあれば、一旦削除
    if(os.path.isdir('fig') == True):
        shutil.rmtree('fig')
    load = convert_from_path(fullPath)
    if not os.path.exists("fig"):
        # figを作る。
        os.makedirs("fig")
    img_file = "fig/output"
    for index, jpg in enumerate(load):
        jpg_name = img_file + '_' + '{0:03d}.jpg'.format(index)
        jpg.save(jpg_name)
        print("変換中_{0:03d}.jpg".format(index))
    print("pdfを画像に変換しました。")
else:
    print("pdfファイルがありません、処理をやめます。")
    sys.exit()


# 画像ファイルの読み込み
fnms = glob('./fig/*.jpg')
print(str(len(fnms)) + "枚の画像を処理します。")

# ソート
fnms.sort(key=str2int)


im = Image.open(fnms[0])
h = im.height
w = im.width

# Presentationインスタンスの作成
ppt = Presentation()

# デフォルトはA4縦にしておく
ppt.slide_height = Pt(h)
ppt.slide_width = Pt(w)

# 幅
width = ppt.slide_width
# 高さ
height = ppt.slide_height

# レイアウト, 6番は白紙
blank_slide_layout = ppt.slide_layouts[6]

# ファイル毎にループ
for fnm in fnms:
    print(fnm + "を処理します。")
    # 白紙のスライドの追加
    slide = ppt.slides.add_slide(blank_slide_layout)

    # 画像の挿入
    pic = slide.shapes.add_picture(fnm, 0, 0)

    # 中心に移動
    pic.left = int((width - pic.width) / 2)
    pic.top = int((height - pic.height) / 2)
    print(fnm + "の処理が終わった！")

# 名前をつけて保存
ppt.save(outputdir +  '/output.pptx')
shutil.rmtree('fig')
print("変換が終わりました！！！ output.pptxを作成しました。")
